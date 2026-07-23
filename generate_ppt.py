"""
Generate a professional PPT introducing the 消防安全智能评估系统 (SafeEvaluate).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour palette ──────────────────────────────────────────────
PRIMARY    = RGBColor(0x1A, 0x56, 0xDB)   # blue
DARK       = RGBColor(0x0F, 0x2B, 0x54)   # dark navy
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG   = RGBColor(0xF8, 0xFA, 0xFC)
TEXT_DARK   = RGBColor(0x1E, 0x29, 0x3B)
TEXT_GRAY   = RGBColor(0x6B, 0x72, 0x80)
GREEN      = RGBColor(0x16, 0xA3, 0x4A)
RED        = RGBColor(0xDC, 0x26, 0x26)
ORANGE     = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE     = RGBColor(0x7C, 0x3A, 0xED)
CYAN       = RGBColor(0x06, 0xB6, 0xD4)
SLIDE_W    = Inches(13.333)
SLIDE_H    = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helper functions ─────────────────────────────────────────────
def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_DARK, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Microsoft YaHei', anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    txBox.text_frame.auto_size = None
    tf = txBox.text_frame
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    try:
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)
    except Exception:
        pass
    return txBox

def add_multiline_box(slide, left, top, width, height, lines, font_name='Microsoft YaHei'):
    """lines is a list of (text, font_size, color, bold, alignment)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    tf = txBox.text_frame
    tf.auto_size = None
    for i, line_data in enumerate(lines):
        text, font_size, color, bold, alignment = line_data
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_before = Pt(2)
        p.space_after = Pt(2)
    return txBox

def add_page_number(slide, num):
    add_text_box(slide, Inches(12.2), Inches(7.0), Inches(0.9), Inches(0.4),
                 str(num), font_size=10, color=TEXT_GRAY, alignment=PP_ALIGN.RIGHT)

def add_section_title(slide, title, subtitle=None):
    """Standard slide header with blue accent bar."""
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.35), Inches(11.7), Inches(0.65),
                 title, font_size=30, color=DARK, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(0.95), Inches(11.7), Inches(0.45),
                     subtitle, font_size=14, color=TEXT_GRAY)
    # Thin separator line
    add_rect(slide, Inches(0.8), Inches(1.35), Inches(11.7), Inches(0.015), RGBColor(0xE5, 0xE7, 0xEB))

# ── Cover slide ──────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK)

# Decorative shapes
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), PRIMARY)
add_rect(slide, Inches(0), Inches(7.38), Inches(13.333), Inches(0.12), PRIMARY)
# Large circle decoration
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-1.5), Inches(6), Inches(6))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x15, 0x3E, 0x75)
circle.line.fill.background()
circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.8), Inches(4.5), Inches(3), Inches(3))
circle2.fill.solid()
circle2.fill.fore_color.rgb = RGBColor(0x15, 0x3E, 0x75)
circle2.line.fill.background()

# Title
add_text_box(slide, Inches(1.2), Inches(1.8), Inches(8), Inches(1.0),
             '消防安全智能评估系统', font_size=48, color=WHITE, bold=True)
add_text_box(slide, Inches(1.2), Inches(2.85), Inches(8), Inches(0.6),
             'SafeEvaluate  |  AI-Powered Fire Safety Assessment', font_size=20, color=RGBColor(0x93, 0xC5, 0xFD))
add_rect(slide, Inches(1.2), Inches(3.55), Inches(2.5), Inches(0.06), PRIMARY)

add_multiline_box(slide, Inches(1.2), Inches(3.9), Inches(7), Inches(1.2), [
    ('智能评估 · 精准研判  |  保障消防安全', 16, RGBColor(0x93, 0xC5, 0xFD), False, PP_ALIGN.LEFT),
    ('', 10, WHITE, False, PP_ALIGN.LEFT),
    ('基于千问视觉大模型 + 消防法规知识库', 14, RGBColor(0x9C, 0xA3, 0xAF), False, PP_ALIGN.LEFT),
])
# Bottom info
add_text_box(slide, Inches(1.2), Inches(6.5), Inches(5), Inches(0.4),
             '公安内部系统 · 仅限授权人员访问', font_size=11, color=TEXT_GRAY)

# ── Slide 2: 项目概述 ────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, '项目概述', '系统定位 · 核心能力 · 技术栈')

# Left: system intro card
card1 = add_rounded_rect(slide, Inches(0.8), Inches(1.7), Inches(3.7), Inches(2.5), LIGHT_BG)
add_text_box(slide, Inches(1.15), Inches(1.9), Inches(3.0), Inches(0.4),
             '🎯  系统定位', font_size=18, color=DARK, bold=True)
add_multiline_box(slide, Inches(1.15), Inches(2.4), Inches(3.2), Inches(1.7), [
    ('面向派出所防火部门的', 13, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('消防安全智能评估平台', 13, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('', 8, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('上传现场照片/图纸，AI 自动', 12, TEXT_GRAY, False, PP_ALIGN.LEFT),
    ('依据消防法规生成评估报告', 12, TEXT_GRAY, False, PP_ALIGN.LEFT),
    ('', 8, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('智能发现隐患 · 辅助执法决策', 12, TEXT_GRAY, False, PP_ALIGN.LEFT),
])

card2 = add_rounded_rect(slide, Inches(4.85), Inches(1.7), Inches(3.7), Inches(2.5), LIGHT_BG)
add_text_box(slide, Inches(5.2), Inches(1.9), Inches(3.0), Inches(0.4),
             '🧠  核心能力', font_size=18, color=DARK, bold=True)
add_multiline_box(slide, Inches(5.2), Inches(2.4), Inches(3.2), Inches(1.7), [
    ('✅ AI 图像识别 — 多图关联分析', 12, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('✅ 法规对标 — 消防标准自动匹配', 12, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('✅ 智能报告 — 一键生成评估结果', 12, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('✅ 规则管理 — 支持自定义检查规则', 12, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('✅ 统计分析 — 合规趋势大屏数据', 12, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('✅ 历史追溯 — 所有报告永久存档', 12, TEXT_DARK, False, PP_ALIGN.LEFT),
])

card3 = add_rounded_rect(slide, Inches(8.9), Inches(1.7), Inches(3.7), Inches(2.5), LIGHT_BG)
add_text_box(slide, Inches(9.25), Inches(1.9), Inches(3.0), Inches(0.4),
             '⚙️  技术栈', font_size=18, color=DARK, bold=True)
add_multiline_box(slide, Inches(9.25), Inches(2.4), Inches(3.2), Inches(1.7), [
    ('前端：React 18 + Vite + CSS Modules', 12, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('后端：Python FastAPI + Uvicorn', 12, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('AI：阿里云千问视觉大模型', 12, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('', 8, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('认证：JWT（HS256 令牌）', 12, TEXT_GRAY, False, PP_ALIGN.LEFT),
    ('存储：JSON 文件持久化', 12, TEXT_GRAY, False, PP_ALIGN.LEFT),
    ('部署：Docker + Nginx 反向代理', 12, TEXT_GRAY, False, PP_ALIGN.LEFT),
])

# Bottom: workflow (simple steps)
add_text_box(slide, Inches(0.8), Inches(4.6), Inches(3), Inches(0.4),
             '📋  评估工作流程', font_size=18, color=DARK, bold=True)

steps = [
    ('1', '上传资料', '图片/PDF/图纸'),
    ('2', '选择规则', '勾选检查标准'),
    ('3', 'AI 分析', '千问视觉大模型'),
    ('4', '生成报告', '合规/隐患明细'),
]
step_w = Inches(2.65)
step_start = Inches(0.8)
for i, (num, title, desc) in enumerate(steps):
    x = step_start + i * (step_w + Inches(0.3))
    # Step circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.9), Inches(5.15), Inches(0.65), Inches(0.65))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PRIMARY
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = 'Microsoft YaHei'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Arrow between (except last)
    if i < 3:
        add_text_box(slide, x + Inches(1.7), Inches(5.2), Inches(1.0), Inches(0.5),
                     '▸', font_size=22, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
    # Labels
    add_text_box(slide, x, Inches(5.9), Inches(2.65), Inches(0.35),
                 title, font_size=14, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(6.25), Inches(2.65), Inches(0.35),
                 desc, font_size=11, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 2)

# ── Slide 3: 登录与权限 ──────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, '🔐  登录与权限管理', 'JWT 认证 · 身份验证 · 安全防护')

# Screenshot placeholder (simulated layout description)
# Left panel
left = add_rounded_rect(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(5.2), DARK)
add_text_box(slide, Inches(1.3), Inches(2.0), Inches(4.5), Inches(0.5),
             '🔑  登录界面', font_size=18, color=WHITE, bold=True)
add_multiline_box(slide, Inches(1.3), Inches(2.6), Inches(4.5), Inches(3.5), [
    ('分屏布局设计', 16, RGBColor(0x93, 0xC5, 0xFD), True, PP_ALIGN.LEFT),
    ('', 8, WHITE, False, PP_ALIGN.LEFT),
    ('左侧品牌面板：', 13, WHITE, True, PP_ALIGN.LEFT),
    ('  · 消防标识 + 系统名称 + 口号', 12, RGBColor(0x9C, 0xA3, 0xAF), False, PP_ALIGN.LEFT),
    ('  · 渐变深蓝底色 + 装饰圆形图案', 12, RGBColor(0x9C, 0xA3, 0xAF), False, PP_ALIGN.LEFT),
    ('', 8, WHITE, False, PP_ALIGN.LEFT),
    ('右侧登录卡片：', 13, WHITE, True, PP_ALIGN.LEFT),
    ('  · 用户名 + 密码输入', 12, RGBColor(0x9C, 0xA3, 0xAF), False, PP_ALIGN.LEFT),
    ('  · 前端表单验证（非空检查）', 12, RGBColor(0x9C, 0xA3, 0xAF), False, PP_ALIGN.LEFT),
    ('  · 登录按钮 + 错误提示', 12, RGBColor(0x9C, 0xA3, 0xAF), False, PP_ALIGN.LEFT),
    ('', 8, WHITE, False, PP_ALIGN.LEFT),
    ('  · 底部标注 "公安内部系统"', 11, RGBColor(0x6B, 0x72, 0x80), False, PP_ALIGN.LEFT),
])

# Right features
features = [
    ('JWT 令牌认证', 'HS256 签名算法，24 小时有效期', PRIMARY),
    ('请求拦截器', 'axios 自动附加 Authorization 头', GREEN),
    ('路由守卫', 'ProtectedRoute 组件包裹受保护页面', ORANGE),
    ('401 自动跳转', '令牌过期自动重定向到登录页', PURPLE),
]
for i, (title, desc, color) in enumerate(features):
    y = Inches(1.7) + i * Inches(1.3)
    add_rounded_rect(slide, Inches(6.7), y, Inches(5.8), Inches(1.1), LIGHT_BG)
    # Color accent dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.95), y + Inches(0.38), Inches(0.3), Inches(0.3))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    add_text_box(slide, Inches(7.45), y + Inches(0.1), Inches(4.5), Inches(0.35),
                 title, font_size=16, color=DARK, bold=True)
    add_text_box(slide, Inches(7.45), y + Inches(0.5), Inches(4.5), Inches(0.35),
                 desc, font_size=12, color=TEXT_GRAY)

add_page_number(slide, 3)

# ── Slide 4: 新建评估（上传 + 规则选择） ──────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, '📤  新建评估 — 文件上传与规则选择', '支持多图上传 · 拖拽操作 · 按类别勾选规则')

# Left column: upload zone
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.4),
             '📷  文件上传区域', font_size=18, color=DARK, bold=True)

upload_card = add_rounded_rect(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(2.6), LIGHT_BG)
# Dashed border simulated
upload_card.line.color.rgb = PRIMARY
upload_card.line.width = Pt(2)

add_text_box(slide, Inches(1.2), Inches(2.5), Inches(4.7), Inches(0.4),
             '拖拽文件到此处  或  点击选择文件', font_size=16, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.2), Inches(3.1), Inches(4.7), Inches(0.4),
             '支持 PNG / JPEG / GIF / BMP / WebP / PDF', font_size=12, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.2), Inches(3.5), Inches(4.7), Inches(0.4),
             '单文件最大 50MB，支持多张图片同时上传', font_size=12, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Upload features list
add_multiline_box(slide, Inches(0.8), Inches(5.0), Inches(5.5), Inches(1.8), [
    ('上传特性：', 14, DARK, True, PP_ALIGN.LEFT),
    ('  · 基于 react-dropzone 实现拖拽上传', 12, TEXT_GRAY, False, PP_ALIGN.LEFT),
    ('  · 已选文件列表：名称 + 大小 + 删除按钮', 12, TEXT_GRAY, False, PP_ALIGN.LEFT),
    ('  · 前端校验文件类型和大小', 12, TEXT_GRAY, False, PP_ALIGN.LEFT),
    ('  · 多张图片合并发送，AI 交叉分析', 12, TEXT_GRAY, False, PP_ALIGN.LEFT),
    ('  · 提交时显示 Loading 动画 (30-90秒)', 12, TEXT_GRAY, False, PP_ALIGN.LEFT),
])

# Right column: rule selector
add_text_box(slide, Inches(6.8), Inches(1.6), Inches(5.5), Inches(0.4),
             '📋  规则选择器（侧边栏）', font_size=18, color=DARK, bold=True)

rule_card = add_rounded_rect(slide, Inches(6.8), Inches(2.1), Inches(5.5), Inches(4.7), LIGHT_BG)
add_multiline_box(slide, Inches(7.2), Inches(2.3), Inches(4.8), Inches(4.3), [
    ('按类别分组的规则列表：', 14, DARK, True, PP_ALIGN.LEFT),
    ('', 6, WHITE, False, PP_ALIGN.LEFT),
    ('🏃  消防通道与疏散', 13, DARK, True, PP_ALIGN.LEFT),
    ('     安全出口宽度 / 疏散通道畅通 / 防火门状态', 11, TEXT_GRAY, False, PP_ALIGN.LEFT),
    ('🧯  消防设施与器材', 13, DARK, True, PP_ALIGN.LEFT),
    ('     灭火器配置 / 自动喷淋系统 / 消火栓状态', 11, TEXT_GRAY, False, PP_ALIGN.LEFT),
    ('⚡  电气与火源管理', 13, DARK, True, PP_ALIGN.LEFT),
    ('     电线敷设 / 漏电保护 / 明火管理', 11, TEXT_GRAY, False, PP_ALIGN.LEFT),
    ('📋  消防安全管理', 13, DARK, True, PP_ALIGN.LEFT),
    ('     巡查记录 / 应急预案 / 值班制度', 11, TEXT_GRAY, False, PP_ALIGN.LEFT),
    ('🏗️  建筑与场所属性', 13, DARK, True, PP_ALIGN.LEFT),
    ('     防火分区 / 耐火等级 / 建筑用途', 11, TEXT_GRAY, False, PP_ALIGN.LEFT),
    ('', 6, WHITE, False, PP_ALIGN.LEFT),
    ('💡 不勾选任何规则 = AI 依据全部法规文档评估', 11, PRIMARY, False, PP_ALIGN.LEFT),
])

add_page_number(slide, 4)

# ── Slide 5: 评估报告 ───────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, '📊  评估报告', 'AI 分析结果展示 · 合规/隐患明细 · 法规依据')

# Stat cards illustration
stats_data = [
    ('过关项', '合规', GREEN, '✅'),
    ('不过关项', '存在隐患', RED, '🔴'),
    ('整改建议', '限期整改', ORANGE, '💡'),
]
for i, (label, desc, color, icon) in enumerate(stats_data):
    x = Inches(0.8) + i * Inches(2.1)
    card = add_rounded_rect(slide, x, Inches(1.55), Inches(1.85), Inches(1.5), WHITE)
    card.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    card.line.width = Pt(1)
    # Top color bar
    add_rect(slide, x, Inches(1.55), Inches(1.85), Inches(0.06), color)
    add_text_box(slide, x + Inches(0.15), Inches(1.75), Inches(1.55), Inches(0.3),
                 icon, font_size=24, color=TEXT_DARK)
    add_text_box(slide, x + Inches(0.15), Inches(2.15), Inches(1.55), Inches(0.3),
                 label, font_size=16, color=DARK, bold=True)
    add_text_box(slide, x + Inches(0.15), Inches(2.45), Inches(1.55), Inches(0.3),
                 desc, font_size=11, color=TEXT_GRAY)

# Left: report structure
add_text_box(slide, Inches(0.8), Inches(3.3), Inches(5.5), Inches(0.4),
             '📋  报告结构', font_size=18, color=DARK, bold=True)
add_multiline_box(slide, Inches(0.8), Inches(3.75), Inches(5.5), Inches(3.2), [
    ('1. 报告标题 + 评估日期', 13, DARK, True, PP_ALIGN.LEFT),
    ('', 6, WHITE, False, PP_ALIGN.LEFT),
    ('2. 总体评估 (Overall Assessment)', 13, DARK, True, PP_ALIGN.LEFT),
    ('   AI 生成的场所整体消防安全状况概述', 11, TEXT_GRAY, False, PP_ALIGN.LEFT),
    ('', 6, WHITE, False, PP_ALIGN.LEFT),
    ('3. 统计卡片', 13, DARK, True, PP_ALIGN.LEFT),
    ('   合规项 / 不合规项 / 整改建议 三大统计', 11, TEXT_GRAY, False, PP_ALIGN.LEFT),
    ('', 6, WHITE, False, PP_ALIGN.LEFT),
    ('4. 不过关项详情 🔴', 13, DARK, True, PP_ALIGN.LEFT),
    ('   隐患标题 + 详细说明 + 法规依据引用', 11, TEXT_GRAY, False, PP_ALIGN.LEFT),
    ('   严重等级标签：danger（不过关）/ warning（建议）', 11, TEXT_GRAY, False, PP_ALIGN.LEFT),
    ('', 6, WHITE, False, PP_ALIGN.LEFT),
    ('5. 过关项详情 🟢', 13, DARK, True, PP_ALIGN.LEFT),
    ('   合规项标题 + 说明 + 对应法规条款', 11, TEXT_GRAY, False, PP_ALIGN.LEFT),
])

# Right: finding card example
add_text_box(slide, Inches(6.8), Inches(3.3), Inches(5.5), Inches(0.4),
             '🔍  评估项卡片示例', font_size=18, color=DARK, bold=True)

finding = add_rounded_rect(slide, Inches(6.8), Inches(3.75), Inches(5.5), Inches(2.8), LIGHT_BG)
# Severity badge
badge = add_rounded_rect(slide, Inches(7.15), Inches(3.95), Inches(1.0), Inches(0.35), RED)
add_text_box(slide, Inches(7.15), Inches(3.95), Inches(1.0), Inches(0.35),
             '不过关', font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Category badge
cat_badge = add_rounded_rect(slide, Inches(8.3), Inches(3.95), Inches(1.6), Inches(0.35), LIGHT_BG)
cat_badge.line.color.rgb = PRIMARY
cat_badge.line.width = Pt(1)
add_text_box(slide, Inches(8.3), Inches(3.95), Inches(1.6), Inches(0.35),
             '消防通道与疏散', font_size=9, color=PRIMARY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(7.15), Inches(4.45), Inches(4.8), Inches(0.35),
             '⚠️ 疏散通道宽度不达标', font_size=15, color=DARK, bold=True)
add_multiline_box(slide, Inches(7.15), Inches(4.85), Inches(4.8), Inches(1.2), [
    ('二层东侧疏散通道实测宽度约1.1m，规范要求疏散', 11, TEXT_GRAY, False, PP_ALIGN.LEFT),
    ('通道净宽度不应小于1.4m，存在人员疏散拥堵风险。', 11, TEXT_GRAY, False, PP_ALIGN.LEFT),
    ('', 6, WHITE, False, PP_ALIGN.LEFT),
    ('法规依据：GB 50016 建筑设计防火规范 第5.5.18条', 11, TEXT_GRAY, False, PP_ALIGN.LEFT),
])

# Action buttons
add_text_box(slide, Inches(6.8), Inches(6.7), Inches(5.5), Inches(0.35),
             '🖨️  打印报告    |    📄  导出 PDF', font_size=13, color=TEXT_GRAY)

add_page_number(slide, 5)

# ── Slide 6: 规则管理 ────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, '📝  规则管理', 'CRUD 操作 · 分类筛选 · 内置规则与自定义规则')

# Left: rule list table
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(6.5), Inches(0.4),
             '📋  规则列表（表格视图）', font_size=18, color=DARK, bold=True)

# Simulate a mini table
table_card = add_rounded_rect(slide, Inches(0.8), Inches(2.1), Inches(6.5), Inches(3.0), WHITE)
table_card.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
table_card.line.width = Pt(1)

# Table header
add_rect(slide, Inches(0.8), Inches(2.1), Inches(6.5), Inches(0.45), DARK)
for i, (text, x) in enumerate([('规则名称', Inches(0.95)), ('类别', Inches(2.65)), ('来源文档', Inches(3.75)), ('类型', Inches(5.45))]):
    add_text_box(slide, x, Inches(2.15), Inches(1.5), Inches(0.32),
                 text, font_size=10, color=WHITE, bold=True)

# Table rows
rows = [
    ('中华人民共和国消防法', '消防安全管理', '消防法', '内置'),
    ('GB 50016 建筑设计防火规范', '建筑与场所属性', 'GB 50016', '内置'),
    ('GB 50116 火灾自动报警规范', '消防设施与器材', 'GB 50116', '内置'),
    ('GB 35181-2025 重大火灾隐患判定', '消防安全管理', 'GB 35181', '内置'),
    ('本场所灭火器检查标准', '消防设施与器材', '自定义', '自定义'),
]
for i, (name, cat, src, typ) in enumerate(rows):
    y = Inches(2.58) + i * Inches(0.42)
    bg_color = LIGHT_BG if i % 2 == 0 else WHITE
    add_rect(slide, Inches(0.8), y, Inches(6.5), Inches(0.42), bg_color)
    add_text_box(slide, Inches(0.95), y + Inches(0.06), Inches(1.6), Inches(0.3),
                 name, font_size=10, color=TEXT_DARK)
    add_text_box(slide, Inches(2.65), y + Inches(0.06), Inches(1.0), Inches(0.3),
                 cat, font_size=9, color=TEXT_GRAY)
    add_text_box(slide, Inches(3.75), y + Inches(0.06), Inches(1.5), Inches(0.3),
                 src, font_size=9, color=TEXT_GRAY)
    type_color = GREEN if typ == '内置' else ORANGE
    add_text_box(slide, Inches(5.45), y + Inches(0.06), Inches(0.85), Inches(0.25),
                 typ, font_size=9, color=type_color, bold=True)

# Right: features description
features_rules = [
    ('🏷️  分类筛选', '6大类别下拉筛选：消防通道与疏散、消防设施与器材、电气与火源管理、消防安全管理、建筑与场所属性、其他。切换类别即时过滤规则列表。', PRIMARY),
    ('➕  新增规则', '弹窗表单：规则名称（必填）、类别选择、描述文本、来源文档、条款编号。保存后即时刷新列表。', GREEN),
    ('✏️  编辑规则', '点击编辑按钮打开预填充表单，修改后保存。内置规则和自定义规则均可编辑。', ORANGE),
    ('🗑️  删除规则', '仅自定义规则可删除，内置规则受保护（返回 403）。删除前弹出确认对话框。', RED),
]
for i, (title, desc, color) in enumerate(features_rules):
    y = Inches(1.6) + i * Inches(1.35)
    # Color left bar
    add_rect(slide, Inches(7.65), y, Inches(0.06), Inches(1.15), color)
    add_rounded_rect(slide, Inches(7.85), y + Inches(0.05), Inches(5.1), Inches(1.1), LIGHT_BG)
    add_text_box(slide, Inches(8.05), y + Inches(0.12), Inches(4.7), Inches(0.3),
                 title, font_size=15, color=DARK, bold=True)
    add_text_box(slide, Inches(8.05), y + Inches(0.45), Inches(4.7), Inches(0.6),
                 desc, font_size=10, color=TEXT_GRAY)

# 8 built-in rules + custom support
add_text_box(slide, Inches(0.8), Inches(5.3), Inches(12), Inches(0.3),
             '📦 系统内置 8 条核心规则（消防法、GB标准、地方界定标准）+ 支持用户自定义添加', font_size=12, color=TEXT_GRAY)

add_page_number(slide, 6)

# ── Slide 7: 历史记录 ────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, '📜  历史记录', '评估报告归档 · 分页浏览 · 风险等级标识')

# Left: history table
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(6.5), Inches(0.4),
             '📋  历史记录列表', font_size=18, color=DARK, bold=True)

hist_card = add_rounded_rect(slide, Inches(0.8), Inches(2.1), Inches(6.5), Inches(3.5), WHITE)
hist_card.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
hist_card.line.width = Pt(1)

# Table header
add_rect(slide, Inches(0.8), Inches(2.1), Inches(6.5), Inches(0.45), DARK)
for text, x in [('评估名称', Inches(0.95)), ('评估日期', Inches(2.65)), ('风险等级', Inches(4.35)), ('操作', Inches(5.65))]:
    add_text_box(slide, x, Inches(2.15), Inches(1.5), Inches(0.32),
                 text, font_size=10, color=WHITE, bold=True)

hist_data = [
    ('万达广场消防评估', '2026-07-20', '低风险', GREEN),
    ('人民医院消防复查', '2026-07-18', '中风险', ORANGE),
    ('阳光小学消防检查', '2026-07-15', '高风险', RED),
    ('写字楼A座消防评估', '2026-07-12', '中风险', ORANGE),
    ('居民小区消防检查', '2026-07-08', '低风险', GREEN),
]
for i, (name, date, risk, color) in enumerate(hist_data):
    y = Inches(2.58) + i * Inches(0.42)
    bg_color = LIGHT_BG if i % 2 == 0 else WHITE
    add_rect(slide, Inches(0.8), y, Inches(6.5), Inches(0.42), bg_color)
    add_text_box(slide, Inches(0.95), y + Inches(0.06), Inches(1.6), Inches(0.3),
                 name, font_size=10, color=TEXT_DARK)
    add_text_box(slide, Inches(2.65), y + Inches(0.06), Inches(1.5), Inches(0.3),
                 date, font_size=10, color=TEXT_GRAY)
    # Risk badge
    badge = add_rounded_rect(slide, Inches(4.35), y + Inches(0.08), Inches(1.0), Inches(0.25), color)
    add_text_box(slide, Inches(4.35), y + Inches(0.08), Inches(1.0), Inches(0.25),
                 risk, font_size=8, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(5.65), y + Inches(0.06), Inches(1.0), Inches(0.3),
                 '查看报告 →', font_size=9, color=PRIMARY)

# Pagination simulation
add_text_box(slide, Inches(3.0), Inches(5.65), Inches(3.0), Inches(0.3),
             '‹  1   2   3  ›', font_size=11, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Right: features
add_text_box(slide, Inches(7.65), Inches(1.6), Inches(5.0), Inches(0.4),
             '✨  功能特性', font_size=18, color=DARK, bold=True)

hist_features = [
    ('📄 分页展示', '每页显示 10 条记录，支持页码切换'),
    ('🏷️ 风险分级', '0项不合规 = 低风险 🟢\n1-3项 = 中风险 🟡\n4+项 = 高风险 🔴'),
    ('🔗 快速跳转', '点击行或"查看报告"链接直接跳转到详细报告页'),
    ('🔄 自动刷新', '路由变化时自动重新加载数据'),
    ('📊 总计数', '页头显示历史记录总条数'),
]
for i, (title, desc) in enumerate(hist_features):
    y = Inches(2.1) + i * Inches(1.0)
    add_rounded_rect(slide, Inches(7.65), y, Inches(5.0), Inches(0.8), LIGHT_BG)
    add_text_box(slide, Inches(7.9), y + Inches(0.1), Inches(4.5), Inches(0.25),
                 title, font_size=13, color=DARK, bold=True)
    add_text_box(slide, Inches(7.9), y + Inches(0.38), Inches(4.5), Inches(0.38),
                 desc, font_size=10, color=TEXT_GRAY)

add_page_number(slide, 7)

# ── Slide 8: 统计分析（上） ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, '📈  统计分析 — 数据总览', '概览卡片 · 风险分布 · 分类分布 · 高频问题')

# Overview cards
overview_data = [
    ('📊', '评估总次数', '128', PRIMARY),
    ('📈', '整体合规率', '72.5%', GREEN),
    ('⚠️', '不合规项总数', '87', RED),
    ('💡', '整改建议总数', '64', ORANGE),
]
for i, (icon, label, value, color) in enumerate(overview_data):
    x = Inches(0.8) + i * Inches(3.1)
    card = add_rounded_rect(slide, x, Inches(1.55), Inches(2.85), Inches(1.3), WHITE)
    card.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    card.line.width = Pt(1)
    add_rect(slide, x, Inches(1.55), Inches(2.85), Inches(0.06), color)
    add_text_box(slide, x + Inches(0.2), Inches(1.7), Inches(0.5), Inches(0.35),
                 icon, font_size=20, color=TEXT_DARK)
    add_text_box(slide, x + Inches(0.2), Inches(2.1), Inches(2.4), Inches(0.35),
                 value, font_size=26, color=DARK, bold=True)
    add_text_box(slide, x + Inches(0.2), Inches(2.5), Inches(2.4), Inches(0.3),
                 label, font_size=11, color=TEXT_GRAY)

# Risk distribution
add_text_box(slide, Inches(0.8), Inches(3.1), Inches(5.5), Inches(0.35),
             '🎯  风险等级分布', font_size=16, color=DARK, bold=True)

risks = [
    ('🟢 低风险', 62, GREEN, 0.62),
    ('🟡 中风险', 38, ORANGE, 0.38),
    ('🔴 高风险', 28, RED, 0.28),
]
for i, (label, count, color, pct) in enumerate(risks):
    y = Inches(3.55) + i * Inches(0.65)
    add_text_box(slide, Inches(1.0), y, Inches(1.5), Inches(0.35),
                 label, font_size=12, color=DARK)
    # Bar track
    add_rounded_rect(slide, Inches(2.6), y + Inches(0.08), Inches(6.5), Inches(0.22), RGBColor(0xE5, 0xE7, 0xEB))
    # Bar fill
    bar_w = int(6.5 * pct)
    if bar_w > 0:
        add_rounded_rect(slide, Inches(2.6), y + Inches(0.08), Inches(bar_w), Inches(0.22), color)
    add_text_box(slide, Inches(9.3), y, Inches(1.0), Inches(0.35),
                 str(count), font_size=12, color=DARK, bold=True)
    add_text_box(slide, Inches(10.0), y, Inches(1.0), Inches(0.35),
                 f'{int(pct*100)}%', font_size=11, color=TEXT_GRAY)

# Top issues + Category distribution
add_text_box(slide, Inches(0.8), Inches(5.35), Inches(5.5), Inches(0.35),
             '📋  不合规项分类分布', font_size=16, color=DARK, bold=True)

cats = [
    ('消防通道与疏散', 28, 32, RGBColor(0xDC, 0x26, 0x26)),
    ('消防设施与器材', 22, 25, RGBColor(0x25, 0x63, 0xEB)),
    ('电气与火源管理', 15, 17, RGBColor(0xD9, 0x77, 0x06)),
    ('消防安全管理', 13, 15, RGBColor(0x7C, 0x3A, 0xED)),
    ('建筑与场所属性', 9, 11, RGBColor(0x05, 0x96, 0x69)),
]
for i, (name, count, pct, color) in enumerate(cats):
    y = Inches(5.75) + i * Inches(0.3)
    add_text_box(slide, Inches(1.0), y, Inches(1.8), Inches(0.25),
                 name, font_size=9, color=TEXT_DARK)
    add_rounded_rect(slide, Inches(2.85), y + Inches(0.05), Inches(3.5), Inches(0.15), RGBColor(0xE5, 0xE7, 0xEB))
    bar_w = int(3.5 * pct / 100)
    if bar_w > 0:
        add_rounded_rect(slide, Inches(2.85), y + Inches(0.05), Inches(bar_w), Inches(0.15), color)
    add_text_box(slide, Inches(6.5), y, Inches(0.6), Inches(0.25),
                 f'{count}项', font_size=9, color=DARK, bold=True)

# Top issues panel
add_text_box(slide, Inches(7.5), Inches(5.35), Inches(5.0), Inches(0.35),
             '🔴  高频问题 TOP 5', font_size=16, color=DARK, bold=True)

issues = [
    (1, '疏散通道宽度不达标', '消防通道与疏散', 12),
    (2, '灭火器数量/位置不符合', '消防设施与器材', 9),
    (3, '安全出口被堵塞或锁闭', '消防通道与疏散', 8),
    (4, '应急照明缺失或损坏', '消防设施与器材', 7),
    (5, '电线乱拉乱接现象', '电气与火源管理', 6),
]
for rank, title, cat, count in issues:
    y = Inches(5.75) + (rank - 1) * Inches(0.3)
    medal = {1: '🥇', 2: '🥈', 3: '🥉'}.get(rank, f'{rank}')
    add_text_box(slide, Inches(7.5), y, Inches(0.4), Inches(0.25),
                 medal, font_size=10, color=TEXT_DARK, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7.95), y, Inches(3.2), Inches(0.25),
                 title, font_size=10, color=TEXT_DARK)
    add_text_box(slide, Inches(11.2), y, Inches(1.0), Inches(0.25),
                 f'{count} 次', font_size=10, color=TEXT_GRAY, bold=True)

add_page_number(slide, 8)

# ── Slide 9: 统计分析（下）— 趋势图 ──────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, '📈  统计分析 — 月度趋势', '合规率趋势 · 评估次数变化 · 数据驱动决策')

# Trend chart simulation
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.4),
             '📊  月度合规率趋势图', font_size=18, color=DARK, bold=True)

chart_card = add_rounded_rect(slide, Inches(0.8), Inches(2.1), Inches(11.7), Inches(4.7), LIGHT_BG)

# Simulated bar chart
months = [
    ('2026-01', 68, 12),
    ('2026-02', 65, 15),
    ('2026-03', 71, 18),
    ('2026-04', 70, 14),
    ('2026-05', 75, 20),
    ('2026-06', 78, 16),
    ('2026-07', 82, 22),
]
bar_area_w = Inches(9.5)
bar_area_h = Inches(3.0)
bar_start_x = Inches(1.8)
bar_start_y = Inches(5.3)
max_bar_h = Inches(2.5)

for i, (month, rate, total) in enumerate(months):
    x = bar_start_x + i * Inches(1.35)
    bar_h = Inches(2.5 * rate / 100)
    y = bar_start_y - bar_h

    # Bar
    bar_color = GREEN if rate >= 75 else ORANGE if rate >= 65 else RED
    add_rect(slide, x, y, Inches(0.7), bar_h, bar_color)

    # Rate label on top
    add_text_box(slide, x - Inches(0.1), y - Inches(0.35), Inches(0.9), Inches(0.3),
                 f'{rate}%', font_size=11, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    # Month label
    add_text_box(slide, x - Inches(0.1), bar_start_y + Inches(0.05), Inches(0.9), Inches(0.25),
                 month, font_size=9, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
    # Count
    add_text_box(slide, x - Inches(0.1), bar_start_y + Inches(0.28), Inches(0.9), Inches(0.2),
                 f'{total}次', font_size=8, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Trend note
add_text_box(slide, Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.35),
             '💡 合规率呈上升趋势，7月达到 82%，反映出消防整改措施的有效落实',
             font_size=12, color=TEXT_GRAY)

add_page_number(slide, 9)

# ── Slide 10: 技术架构 ──────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, '🏗️  技术架构', '前后端分离 · AI 驱动 · 容器化部署')

# Architecture layers
layers = [
    ('🖥️  前端层', 'React 18 + Vite 5\nReact Router 6 + CSS Modules\nreact-dropzone 文件上传\naxios HTTP 客户端', PRIMARY, Inches(2.8)),
    ('🔗  网关层', 'Nginx Alpine\n反向代理 / 静态资源\nAPI 路由转发\n请求限流', RGBColor(0x06, 0xB6, 0xD4), Inches(2.8)),
    ('⚙️  后端层', 'FastAPI (Python 3.10)\nPydantic 数据校验\nPyJWT 身份认证\npython-docx 文档解析', PURPLE, Inches(2.8)),
    ('🧠  AI 层', '阿里云 DashScope API\nQwen3.7-vl-plus 视觉模型\nBase64 图像编码传输\nJSON Schema 结构化输出', ORANGE, Inches(2.8)),
    ('💾  数据层', 'JSON 文件持久化\n报告存储 (reports/)\n规则存储 (rules.json)\n法规文档库 (requirement/)', GREEN, Inches(2.8)),
]

for i, (title, desc, color, height) in enumerate(layers):
    x = Inches(0.8) + i * Inches(2.45)
    card = add_rounded_rect(slide, x, Inches(1.7), Inches(2.25), Inches(3.2), WHITE)
    card.line.color.rgb = color
    card.line.width = Pt(2)
    # Top accent bar
    add_rect(slide, x, Inches(1.7), Inches(2.25), Inches(0.07), color)
    add_text_box(slide, x + Inches(0.15), Inches(1.9), Inches(1.95), Inches(0.35),
                 title, font_size=14, color=DARK, bold=True)
    add_text_box(slide, x + Inches(0.15), Inches(2.35), Inches(1.95), Inches(2.2),
                 desc, font_size=11, color=TEXT_GRAY)

# Bottom: data flow
add_text_box(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.35),
             '🔄  数据流转', font_size=16, color=DARK, bold=True)

flow_steps = [
    ('用户上传', '图片/PDF'),
    ('→', ''),
    ('FastAPI\n接收校验', ''),
    ('→', ''),
    ('Base64\n编码图像', ''),
    ('→', ''),
    ('Qwen VL\n视觉分析', ''),
    ('→', ''),
    ('JSON\n解析结果', ''),
    ('→', ''),
    ('生成报告\n持久化存储', ''),
    ('→', ''),
    ('前端渲染\n展示报告', ''),
]
for i, (title, sub) in enumerate(flow_steps):
    x = Inches(0.4) + i * Inches(1.05)
    if title == '→':
        add_text_box(slide, x, Inches(5.75), Inches(0.6), Inches(0.35),
                     '→', font_size=18, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
    else:
        add_rounded_rect(slide, x, Inches(5.6), Inches(0.95), Inches(0.75), LIGHT_BG)
        add_text_box(slide, x + Inches(0.05), Inches(5.65), Inches(0.85), Inches(0.6),
                     title, font_size=9, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
        if sub:
            add_text_box(slide, x + Inches(0.05), Inches(6.05), Inches(0.85), Inches(0.25),
                         sub, font_size=7, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Docker deploy info
add_text_box(slide, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.4),
             '🐳  部署：Docker + Docker Compose（4 Workers × Uvicorn + Nginx 反向代理）  |  一键部署脚本 + 环境变量配置',
             font_size=11, color=TEXT_GRAY)

add_page_number(slide, 10)

# ── Slide 11: 功能总览全景 ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, '🌟  功能全景图', '六大核心模块 · 覆盖消防安全评估全流程')

modules = [
    ('🔐', '登录认证', 'JWT 令牌\n会话管理\n权限控制', PRIMARY),
    ('📤', '新建评估', '图片上传\n规则选择\nAI 分析', GREEN),
    ('📊', '评估报告', '统计卡片\n隐患明细\n法规引用', ORANGE),
    ('📝', '规则管理', 'CRUD 操作\n分类筛选\n自定义规则', PURPLE),
    ('📜', '历史记录', '分页浏览\n风险分级\n快速检索', CYAN),
    ('📈', '统计分析', '概览卡片\n风险分布\n趋势图表', RED),
]

for i, (icon, title, desc, color) in enumerate(modules):
    x = Inches(0.6) + i * Inches(2.05)
    card = add_rounded_rect(slide, x, Inches(1.7), Inches(1.85), Inches(2.5), WHITE)
    card.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    card.line.width = Pt(1)
    # Icon circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.55), Inches(1.9), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text_box(slide, x + Inches(0.55), Inches(2.0), Inches(0.7), Inches(0.5),
                 icon, font_size=22, color=WHITE, alignment=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide, x + Inches(0.1), Inches(2.75), Inches(1.65), Inches(0.35),
                 title, font_size=15, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    # Description
    add_text_box(slide, x + Inches(0.1), Inches(3.15), Inches(1.65), Inches(0.85),
                 desc, font_size=10, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom highlights
add_text_box(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(0.35),
             '✨  系统亮点', font_size=18, color=DARK, bold=True)

highlights = [
    ('🤖 AI 驱动', '千问视觉大模型自动分析\n消防现场图片，精准识别隐患'),
    ('📋 法规对标', '内置消防法规知识库\n自动匹配评估依据条款'),
    ('🔧 灵活配置', '支持自定义检查规则\n适应不同场所评估需求'),
    ('📊 数据洞察', '多维度统计分析\n助力消防管理决策'),
]
for i, (title, desc) in enumerate(highlights):
    x = Inches(0.8) + i * Inches(3.1)
    add_rounded_rect(slide, x, Inches(5.05), Inches(2.85), Inches(1.9), LIGHT_BG)
    add_text_box(slide, x + Inches(0.15), Inches(5.2), Inches(2.55), Inches(0.35),
                 title, font_size=15, color=DARK, bold=True)
    add_text_box(slide, x + Inches(0.15), Inches(5.65), Inches(2.55), Inches(1.1),
                 desc, font_size=11, color=TEXT_GRAY)

add_page_number(slide, 11)

# ── Slide 12: 感谢页 ────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), PRIMARY)
add_rect(slide, Inches(0), Inches(7.38), Inches(13.333), Inches(0.12), PRIMARY)

# Decoration
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(4.5), Inches(5), Inches(5))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x15, 0x3E, 0x75)
circle.line.fill.background()

add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.0),
             '感谢观看', font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.6),
             '消防安全智能评估系统  |  SafeEvaluate', font_size=22, color=RGBColor(0x93, 0xC5, 0xFD), alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(5.4), Inches(4.0), Inches(2.5), Inches(0.05), PRIMARY)

add_multiline_box(slide, Inches(1.5), Inches(4.3), Inches(10), Inches(1.5), [
    ('智能评估 · 精准研判', 18, RGBColor(0x9C, 0xA3, 0xAF), False, PP_ALIGN.CENTER),
    ('', 10, WHITE, False, PP_ALIGN.CENTER),
    ('保障消防安全，AI 赋能执法', 16, RGBColor(0x9C, 0xA3, 0xAF), False, PP_ALIGN.CENTER),
])

add_text_box(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.4),
             '公安内部系统 · 仅限授权人员访问', font_size=11, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 12)

# ── Save ─────────────────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'SafeEvaluate_产品介绍.pptx')
prs.save(output_path)
print(f'PPT saved to: {output_path}')
print(f'Total slides: {len(prs.slides)}')
